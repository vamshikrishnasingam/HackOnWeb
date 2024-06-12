from pptx import Presentation

def allowed_file(filename, allowed_extensions):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    j=0
    for i in text:
        print(i)
    return " ".join(text)
